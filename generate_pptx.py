from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Premium Palette ──
NAVY = RGBColor(0x0F, 0x1A, 0x2E)
DEEP_NAVY = RGBColor(0x0B, 0x13, 0x24)
GOLD = RGBColor(0xD4, 0xA8, 0x3C)
LIGHT_GOLD = RGBColor(0xF5, 0xE6, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF8, 0xFA)
CARD_BG = RGBColor(0xF0, 0xF2, 0xF6)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
BODY_TEXT = RGBColor(0x44, 0x44, 0x5A)
MUTED = RGBColor(0x88, 0x88, 0x9A)
SUBTLE_BORDER = RGBColor(0xE2, 0xE4, 0xEA)
GREEN = RGBColor(0x0E, 0xA5, 0x6E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
BLUE_ACCENT = RGBColor(0x3B, 0x82, 0xF6)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_flat_rect(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, w, h, lines, size=14, color=BODY_TEXT, bold=False, spacing=6):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(spacing)
    return txBox

def add_table_styled(slide, left, top, w, h, headers, rows, col_widths=None):
    tbl = slide.shapes.add_table(len(rows)+1, len(headers), left, top, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs:
            p.font.color.rgb = GOLD
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci)
            c.text = val
            bg = OFF_WHITE if ri % 2 == 0 else WHITE
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            for p in c.text_frame.paragraphs:
                p.font.color.rgb = DARK_TEXT
                p.font.size = Pt(11)
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER
    return tbl

def add_badge(slide, left, top, text, bg_color, text_color=WHITE, size=10):
    s = add_rect(slide, left, top, Inches(1.3), Inches(0.35), bg_color, 0.15)
    add_text(slide, left+Inches(0.05), top+Inches(0.03), Inches(1.2), Inches(0.3), text, size, text_color, True, PP_ALIGN.CENTER)
    return s

def add_card(slide, left, top, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = SUBTLE_BORDER
    shape.line.width = Pt(0.5)
    shape.adjustments[0] = 0.05
    return shape

# ═══════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DEEP_NAVY)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GOLD)
add_flat_rect(sl, Inches(6.2), Inches(3.15), Inches(0.06), Inches(1.2), GOLD)
add_text(sl, Inches(1), Inches(1.8), Inches(11), Inches(1.2), 'KHQR PAYMENT INTEGRATION', 44, WHITE, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.6), Inches(11), Inches(0.7), 'Bakong (NBC) + ABA Bank', 24, GOLD, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(4.3), Inches(11), Inches(0.5), 'Integration Guide for NexusFinance', 18, MUTED, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.5), 'National Bank of Cambodia  |  ABA Bank  |  KHQR Standard', 12, RGBColor(0x55, 0x66, 0x7A), False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_text(sl, Inches(0.8), Inches(0.4), Inches(4), Inches(0.6), 'AGENDA', 28, NAVY, True, PP_ALIGN.LEFT)
add_flat_rect(sl, Inches(0.8), Inches(1), Inches(0.8), Inches(0.04), GOLD)

sections = [
    ('01', 'Bakong (NBC) KHQR Integration', 'Overview, prerequisites, API reference, KHQR format'),
    ('02', 'ABA Bank API Integration', 'Differences, prerequisites, how to apply, endpoints'),
    ('03', 'Current Implementation Status', 'What exists (mock vs real), what needs to change'),
    ('04', 'Environment Variables', 'Bakong & ABA .env configuration'),
    ('05', 'Quick Start Checklist', 'Step-by-step from registration to go-live'),
]
for i, (num, title, desc) in enumerate(sections):
    y = Inches(1.6) + i * Inches(1.1)
    add_rect(sl, Inches(0.8), y, Inches(0.65), Inches(0.65), GOLD, 0.15)
    add_text(sl, Inches(0.8), y+Inches(0.08), Inches(0.65), Inches(0.5), num, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text(sl, Inches(1.7), y+Inches(0.05), Inches(9), Inches(0.35), title, 16, NAVY, True)
    add_text(sl, Inches(1.7), y+Inches(0.38), Inches(9), Inches(0.3), desc, 12, MUTED)

# ═══════════════════════════════════════════════
# SLIDE 3: BAKONG OVERVIEW
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_badge(sl, Inches(0.8), Inches(0.4), 'PART 1', NAVY, GOLD)
add_text(sl, Inches(2.3), Inches(0.35), Inches(8), Inches(0.5), 'Bakong (NBC) KHQR Integration', 26, NAVY, True)
add_text(sl, Inches(2.3), Inches(0.85), Inches(8), Inches(0.3), "Cambodia's national real-time payment system", 13, MUTED)

# How it works card
add_card(sl, Inches(0.8), Inches(1.5), Inches(5.8), Inches(3.2))
add_text(sl, Inches(1.1), Inches(1.7), Inches(3), Inches(0.4), 'How It Works', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(2.1), Inches(0.5), Inches(0.03), GOLD)
steps = ['Customer scans KHQR code', 'Opens Bakong or any banking app', 'Confirms payment amount', 'NBC clears funds in real-time', 'Webhook notifies your server']
for i, step in enumerate(steps):
    y = Inches(2.3) + i * Inches(0.42)
    add_rect(sl, Inches(1.2), y+Inches(0.05), Inches(0.2), Inches(0.2), GOLD if i < 5 else GREEN, 0.5)
    add_text(sl, Inches(1.6), y, Inches(5), Inches(0.3), step, 12, DARK_TEXT)

# Prerequisites card
add_card(sl, Inches(7), Inches(1.5), Inches(5.5), Inches(3.2))
add_text(sl, Inches(7.3), Inches(1.7), Inches(4), Inches(0.4), 'Prerequisites', 16, NAVY, True)
add_flat_rect(sl, Inches(7.3), Inches(2.1), Inches(0.5), Inches(0.03), GOLD)
add_table_styled(sl, Inches(7.2), Inches(2.4), Inches(5.1), Inches(2),
    ['Item', 'How to Get'],
    [['Bakong Wallet', 'Download Bakong app, register'],
     ['Bakong Account ID', 'Format: yourname@bank (e.g. merchant@aclb)'],
     ['NBC API Token', 'Register with NBC developer support'],
     ['Server in Cambodia', 'Required for production check-transaction']],
    [Inches(2), Inches(3.1)])

# Key info card
add_card(sl, Inches(0.8), Inches(5), Inches(11.7), Inches(1.8))
add_text(sl, Inches(1.1), Inches(5.2), Inches(3), Inches(0.4), 'How to Get NBC API Token', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(5.6), Inches(0.5), Inches(0.03), GOLD)
add_multiline(sl, Inches(1.1), Inches(5.8), Inches(11), Inches(1),
    ['1.  Email NBC developer support or contact your bank partner',
     '2.  They provide a username/password pair for the API',
     '3.  Call POST /v1/request_token to receive a Bearer token',
     '4.  Token expires — use /v1/renew_token to refresh'],
    12, BODY_TEXT, False, 4)

# ═══════════════════════════════════════════════
# SLIDE 4: BAKONG API
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_text(sl, Inches(0.8), Inches(0.4), Inches(6), Inches(0.5), 'Bakong API Reference', 26, NAVY, True)
add_text(sl, Inches(0.8), Inches(0.85), Inches(11), Inches(0.3), 'Base URLs', 13, MUTED)

add_table_styled(sl, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.7),
    ['Environment', 'URL'],
    [['Sandbox (SIT)', 'https://sit-api-bakong.nbc.gov.kh/v1'],
     ['Production', 'https://api-bakong.nbc.gov.kh/v1']],
    [Inches(2), Inches(9.7)])

add_text(sl, Inches(0.8), Inches(2.2), Inches(4), Inches(0.4), 'API Endpoints', 16, NAVY, True)
add_flat_rect(sl, Inches(0.8), Inches(2.6), Inches(0.5), Inches(0.03), GOLD)
add_table_styled(sl, Inches(0.8), Inches(2.8), Inches(11.7), Inches(3),
    ['Method', 'Endpoint', 'Description'],
    [['POST', '/v1/request_token', 'Get access token with credentials'],
     ['POST', '/v1/verify', 'Verify token is still valid'],
     ['POST', '/v1/renew_token', 'Refresh an expiring token'],
     ['POST', '/v1/generate_deeplink_by_qr', 'Generate payment deeplink from KHQR'],
     ['POST', '/v1/check_transaction_by_md5', 'Check payment status by MD5 hash'],
     ['POST', '/v1/check_transaction_by_hash', 'Check payment by full transaction hash'],
     ['POST', '/v1/check_transaction_by_short_hash', 'Check payment by short hash'],
     ['POST', '/v1/check_bakong_account', 'Verify if a Bakong account exists']],
    [Inches(1.2), Inches(5.5), Inches(5)])

add_text(sl, Inches(0.8), Inches(6), Inches(2), Inches(0.3), 'Response Codes:', 11, MUTED)
add_table_styled(sl, Inches(2.2), Inches(6), Inches(4), Inches(0.6),
    ['Code', 'Meaning'],
    [['0', 'Success'], ['1', 'Fail / Not found'], ['6', 'Unauthorized']],
    [Inches(0.8), Inches(3.2)])

# ═══════════════════════════════════════════════
# SLIDE 5: KHQR FORMAT
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5), 'KHQR String Format (EMVCo)', 26, NAVY, True)
add_text(sl, Inches(0.8), Inches(0.85), Inches(8), Inches(0.3), 'KHQR follows the EMVCo QR standard with TLV (Tag-Length-Value) data', 13, MUTED)

add_table_styled(sl, Inches(0.8), Inches(1.4), Inches(11.7), Inches(4.5),
    ['Tag', 'Field', 'Description', 'Example'],
    [['00', 'Payload Format Indicator', 'EMVCo fixed value', '01'],
     ['01', 'Point of Initiation Method', '11 = Static QR, 12 = Dynamic QR', '12'],
     ['29', 'Bakong Account ID', 'Merchant account identifier', 'merchant@aclb'],
     ['53', 'Transaction Currency', '840 = USD, 116 = KHR', '840'],
     ['54', 'Transaction Amount', 'Optional for dynamic QR', '25.00'],
     ['58', 'Country Code', 'ISO 3166-1 alpha-2', 'KH'],
     ['59', 'Merchant Name', 'Display name shown to customer', 'NexusFinance'],
     ['60', 'Merchant City', 'City of business', 'Phnom Penh'],
     ['61', 'Store Label', 'Optional branch or store ID', 'Main Branch'],
     ['63', 'CRC', 'CRC16-CCITT checksum (last 4 chars)', '7D00']],
    [Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.9)])

add_text(sl, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4), 'Example KHQR String:  00020101021229240020merchant@aclb520459995303116540425...63047D00', 10, MUTED)

# ═══════════════════════════════════════════════
# SLIDE 6: ABA OVERVIEW
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_badge(sl, Inches(0.8), Inches(0.4), 'PART 2', NAVY, GOLD)
add_text(sl, Inches(2.3), Inches(0.35), Inches(8), Inches(0.5), 'ABA Bank API Integration', 26, NAVY, True)

# Comparison card
add_card(sl, Inches(0.8), Inches(1.3), Inches(5.8), Inches(3))
add_text(sl, Inches(1.1), Inches(1.5), Inches(4), Inches(0.4), 'Key Differences from Bakong', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(1.9), Inches(0.5), Inches(0.03), GOLD)
add_table_styled(sl, Inches(1.1), Inches(2.2), Inches(5.3), Inches(1.8),
    ['Feature', 'Bakong', 'ABA'],
    [['Coverage', 'All Cambodian banks', 'ABA accounts only'],
     ['QR Standard', 'KHQR (EMVCo)', 'ABA-specific QR'],
     ['API Access', 'Apply via NBC', 'Apply via ABA branch'],
     ['Settlement', 'Real-time via NBC', 'Batch settlement']],
    [Inches(1.3), Inches(2), Inches(2)])

# How to apply card
add_card(sl, Inches(7), Inches(1.3), Inches(5.5), Inches(3))
add_text(sl, Inches(7.3), Inches(1.5), Inches(4), Inches(0.4), 'How to Apply', 16, NAVY, True)
add_flat_rect(sl, Inches(7.3), Inches(1.9), Inches(0.5), Inches(0.03), GOLD)
add_multiline(sl, Inches(7.3), Inches(2.2), Inches(4.8), Inches(2.2),
    ['1. Visit any ABA Bank branch with:',
     '     Business registration / company license',
     '     ID documents',
     '     Bank account statement',
     '',
     '2. Request "Merchant Payment Gateway"',
     '',
     '3. ABA provides: Merchant ID, API Key &',
     '     Secret, Endpoint URLs, IP whitelist'],
    11, BODY_TEXT, False, 3)

# Prerequisites card
add_card(sl, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2))
add_text(sl, Inches(1.1), Inches(4.8), Inches(4), Inches(0.4), 'Prerequisites', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(5.2), Inches(0.5), Inches(0.03), GOLD)
add_table_styled(sl, Inches(1.1), Inches(5.5), Inches(11), Inches(0.9),
    ['Item', 'How to Get'],
    [['ABA Merchant Account', 'Apply at any ABA branch with business documents'],
     ['API Credentials', 'Contact ABA Business Banking for API Key & Secret'],
     ['IP Whitelist', 'Register your server IP address with ABA']],
    [Inches(2), Inches(9)])

# ═══════════════════════════════════════════════
# SLIDE 7: ABA API & FLOW
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_text(sl, Inches(0.8), Inches(0.4), Inches(6), Inches(0.5), 'ABA Bank API Endpoints & Flow', 26, NAVY, True)

# Endpoints card
add_card(sl, Inches(0.8), Inches(1.2), Inches(5.8), Inches(2.2))
add_text(sl, Inches(1.1), Inches(1.4), Inches(4), Inches(0.4), 'Common Endpoints', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(1.8), Inches(0.5), Inches(0.03), GOLD)
add_table_styled(sl, Inches(1.1), Inches(2.1), Inches(5.3), Inches(1.2),
    ['Method', 'Endpoint', 'Description'],
    [['POST', '/api/v1/payment/init', 'Create payment request'],
     ['POST', '/api/v1/payment/status', 'Check payment status'],
     ['POST', '/api/v1/payment/refund', 'Process refund'],
     ['POST', '/api/v1/account/balance', 'Check balance']],
    [Inches(1), Inches(2.5), Inches(1.8)])

# Payment flow card
add_card(sl, Inches(7), Inches(1.2), Inches(5.5), Inches(5.2))
add_text(sl, Inches(7.3), Inches(1.4), Inches(4), Inches(0.4), 'Payment Flow', 16, NAVY, True)
add_flat_rect(sl, Inches(7.3), Inches(1.8), Inches(0.5), Inches(0.03), GOLD)
flow = [
    ('1', 'Initiate', 'Frontend calls your API to start payment'),
    ('2', 'Request', 'You call ABA /api/v1/payment/init'),
    ('3', 'Receive', 'ABA returns payment URL + QR code'),
    ('4', 'Pay', 'Customer pays via ABA Mobile or QR scan'),
    ('5', 'Callback', 'ABA sends webhook to your callbackUrl'),
    ('6', 'Confirm', 'You verify signature, update order status'),
]
for i, (num, title, desc) in enumerate(flow):
    y = Inches(2.1) + i * Inches(0.65)
    add_rect(sl, Inches(7.4), y, Inches(0.3), Inches(0.3), GOLD, 0.5)
    add_text(sl, Inches(7.4), y-Inches(0.02), Inches(0.3), Inches(0.3), num, 9, WHITE, True, PP_ALIGN.CENTER)
    add_text(sl, Inches(7.9), y-Inches(0.02), Inches(1.2), Inches(0.3), title, 11, DARK_TEXT, True)
    add_text(sl, Inches(9.1), y-Inches(0.02), Inches(3), Inches(0.3), desc, 10, MUTED)

# Sample request card
add_card(sl, Inches(0.8), Inches(3.7), Inches(5.8), Inches(2.7))
add_text(sl, Inches(1.1), Inches(3.9), Inches(4), Inches(0.4), 'Initiate Payment — Sample Request', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(4.3), Inches(0.5), Inches(0.03), GOLD)
code = 'POST /api/v1/payment/init\n\nHeaders:\n  X-Merchant-Id: {{merchantId}}\n  X-Signature: {{hmac_sha256}}\n  X-Timestamp: 2025-01-15T10:00:00Z\n\nBody:\n{\n  "txnId": "ORDER_001",\n  "amount": 25.00,\n  "currency": "USD",\n  "callbackUrl": "https://.../api/payment/callback"\n}'
txBox = sl.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(5.2), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code
p.font.name = 'Consolas'
p.font.size = Pt(9)
p.font.color.rgb = BODY_TEXT

# Callback card
add_card(sl, Inches(0.8), Inches(6.7), Inches(5.8), Inches(0.7))
cb_text = 'Callback Webhook: ABA POSTs to your URL with txnId, status, amount, abaTxnRef, signature'
add_text(sl, Inches(1.1), Inches(6.8), Inches(5.3), Inches(0.5), cb_text, 10, BLUE_ACCENT)

# ═══════════════════════════════════════════════
# SLIDE 8: CURRENT STATUS
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5), 'Current Implementation Status', 26, NAVY, True)
add_text(sl, Inches(0.8), Inches(0.85), Inches(8), Inches(0.3), 'What exists now vs what needs to change for production', 13, MUTED)

# What exists
add_card(sl, Inches(0.8), Inches(1.4), Inches(5.8), Inches(2.5))
add_text(sl, Inches(1.1), Inches(1.6), Inches(4), Inches(0.4), 'What Currently Exists', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(2), Inches(0.5), Inches(0.03), GOLD)
add_table_styled(sl, Inches(1.1), Inches(2.3), Inches(5.3), Inches(1.5),
    ['File', 'Status', 'Description'],
    [['server/khqr.ts', 'MOCK', 'EMVCo QR generation + verify'],
     ['KHQRPage.tsx', 'UI DONE', '5-tab frontend (Generate, Verify, Decode, Deeplink, Check)'],
     ['server/index.ts', 'ROUTES', '/api/khqr/* endpoints wired']],
    [Inches(1.5), Inches(1.2), Inches(2.6)])

# Mock vs real
add_card(sl, Inches(7), Inches(1.4), Inches(5.5), Inches(2.5))
add_text(sl, Inches(7.3), Inches(1.6), Inches(4), Inches(0.4), 'Mock vs Real', 16, NAVY, True)
add_flat_rect(sl, Inches(7.3), Inches(2), Inches(0.5), Inches(0.03), GOLD)
add_table_styled(sl, Inches(7.3), Inches(2.3), Inches(5), Inches(1.5),
    ['Feature', 'Current (Mock)', 'Real'],
    [['QR Generation', 'Local EMVCo builder', 'NBC KHQR SDK'],
     ['Payment Check', 'Random (30% success)', 'Bakong API call'],
     ['Deeplink', 'Fake bakong:// link', 'Bakong generate_deeplink'],
     ['Storage', 'In-memory Map', 'Database'],
     ['Bank Coverage', 'Bakong only', 'Bakong + ABA']],
    [Inches(1.4), Inches(2), Inches(1.6)])

# What needs to change
add_card(sl, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5))
add_text(sl, Inches(1.1), Inches(4.5), Inches(6), Inches(0.4), 'What Needs to Change for Production', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(4.9), Inches(0.5), Inches(0.03), GOLD)

changes = [
    'Replace mock functions with real Bakong API calls in server/khqr.ts',
    'Add server/aba.ts module for ABA Bank payment flow',
    'Store transactions in database instead of in-memory Map',
    'Add POST /api/payment/callback webhook endpoint',
    'Add BAKONG_* and ABA_API_* environment variables to .env',
    'Deploy from Cambodia or use proxy for production Bakong API',
]
for i, ch in enumerate(changes):
    y = Inches(5.2) + i * Inches(0.32)
    add_rect(sl, Inches(1.3), y+Inches(0.05), Inches(0.15), Inches(0.15), GOLD, 0.5)
    add_text(sl, Inches(1.7), y, Inches(10.5), Inches(0.3), ch, 11, BODY_TEXT)

# ═══════════════════════════════════════════════
# SLIDE 9: ENV VARS
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_text(sl, Inches(0.8), Inches(0.4), Inches(6), Inches(0.5), 'Environment Variables', 26, NAVY, True)

# Bakong vars
add_card(sl, Inches(0.8), Inches(1.2), Inches(5.8), Inches(2.5))
add_text(sl, Inches(1.1), Inches(1.4), Inches(4), Inches(0.4), 'Bakong (NBC)', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(1.8), Inches(0.5), Inches(0.03), GOLD)
code = 'BAKONG_API_BASE_URL\n    https://sit-api-bakong.nbc.gov.kh/v1\nBAKONG_USERNAME\n    your_email@example.com\nBAKONG_PASSWORD\n    your_api_password\nBAKONG_MERCHANT_ID\n    merchant@aclb\nBAKONG_MERCHANT_NAME\n    NexusFinance'
txBox = sl.shapes.add_textbox(Inches(1.1), Inches(2), Inches(5.3), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code
p.font.name = 'Consolas'
p.font.size = Pt(9.5)
p.font.color.rgb = BODY_TEXT

# ABA vars
add_card(sl, Inches(7), Inches(1.2), Inches(5.5), Inches(2.5))
add_text(sl, Inches(7.3), Inches(1.4), Inches(4), Inches(0.4), 'ABA Bank', 16, NAVY, True)
add_flat_rect(sl, Inches(7.3), Inches(1.8), Inches(0.5), Inches(0.03), GOLD)
code2 = 'ABA_API_BASE_URL\n    https://api.aba.com.kh/v1\nABA_MERCHANT_ID\n    your_aba_merchant_id\nABA_API_KEY\n    your_aba_api_key\nABA_API_SECRET\n    your_aba_api_secret\nABA_CALLBACK_URL\n    https://ndx.com/api/payment/callback'
txBox = sl.shapes.add_textbox(Inches(7.3), Inches(2), Inches(5), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code2
p.font.name = 'Consolas'
p.font.size = Pt(9.5)
p.font.color.rgb = BODY_TEXT

# Full env sample
add_card(sl, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.5))
add_text(sl, Inches(1.1), Inches(4.3), Inches(6), Inches(0.4), 'Full .env Configuration', 16, NAVY, True)
add_flat_rect(sl, Inches(1.1), Inches(4.7), Inches(0.5), Inches(0.03), GOLD)
code3 = '# ── Bakong (NBC) ─────────────────────────────\nBAKONG_API_BASE_URL=https://sit-api-bakong.nbc.gov.kh/v1\nBAKONG_USERNAME=your_email@example.com\nBAKONG_PASSWORD=your_api_password\nBAKONG_MERCHANT_ID=merchant@aclb\nBAKONG_MERCHANT_NAME=NexusFinance\n\n# ── ABA Bank ───────────────────────────────────\nABA_API_BASE_URL=https://api.aba.com.kh/v1\nABA_MERCHANT_ID=your_merchant_id\nABA_API_KEY=your_api_key\nABA_API_SECRET=your_api_secret\nABA_CALLBACK_URL=https://nexusfinancefintech.vercel.app/api/payment/callback'
txBox = sl.shapes.add_textbox(Inches(1.1), Inches(4.9), Inches(11), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = code3
p.font.name = 'Consolas'
p.font.size = Pt(9)
p.font.color.rgb = DARK_TEXT

# ═══════════════════════════════════════════════
# SLIDE 10: CHECKLIST
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, OFF_WHITE)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
add_text(sl, Inches(0.8), Inches(0.4), Inches(6), Inches(0.5), 'Quick Start Checklist', 26, NAVY, True)
add_text(sl, Inches(0.8), Inches(0.85), Inches(8), Inches(0.3), 'Steps to go from mock to live payment integration', 13, MUTED)

items = [
    ('Bakong Setup', 'Download Bakong app, register wallet, get Account ID'),
    ('NBC Token', 'Email NBC or bank for API token'),
    ('Configure Bakong', 'Set BAKONG_API_* vars in .env'),
    ('Update khqr.ts', 'Replace mock with real Bakong API calls'),
    ('ABA Application', 'Visit ABA branch with business docs'),
    ('Configure ABA', 'Set ABA_API_* vars in .env'),
    ('ABA IP Whitelist', 'Register server IP with ABA'),
    ('Create aba.ts', 'Build ABA module for payment flow'),
    ('Webhook', 'Add /api/payment/callback endpoint'),
    ('Sandbox Test', 'Test both Bakong and ABA flows'),
    ('Go Live', 'Deploy to production'),
]

for i, (title, desc) in enumerate(items):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.4) + row * Inches(0.55)
    bg_color = NAVY if i < 3 else (GOLD if i < 6 else CARD_BG)
    text_color = WHITE if i < 6 else DARK_TEXT
    sub_color = LIGHT_GOLD if i < 3 else (LIGHT_GOLD if i < 6 else MUTED)
    add_rect(sl, x, y, Inches(5.8), Inches(0.45), bg_color, 0.05)
    add_text(sl, x+Inches(0.2), y+Inches(0.02), Inches(1.5), Inches(0.2), f'{i+1:02d}. {title}', 10, text_color, True)
    add_text(sl, x+Inches(1.8), y+Inches(0.02), Inches(3.8), Inches(0.2), desc, 9, sub_color)

# ═══════════════════════════════════════════════
# SLIDE 11: THANK YOU
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, DEEP_NAVY)
add_flat_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GOLD)
add_rect(sl, Inches(5.8), Inches(3.3), Inches(1.7), Inches(0.06), GOLD, 0.5)
add_text(sl, Inches(1), Inches(2), Inches(11.333), Inches(1.2), 'Thank You', 48, WHITE, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.7), Inches(11.333), Inches(0.6), 'NexusFinance — Smart Lending Platform', 22, GOLD, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(5), Inches(11.333), Inches(0.5), 'https://nexusfinancefintech.vercel.app', 14, MUTED, False, PP_ALIGN.CENTER)

# ── Save ──
prs.save('C:\\Users\\Asus\\Desktop\\yoooo\\KHQR_Integration_Guide_v2.pptx')
print('Saved: KHQR_Integration_Guide_v2.pptx')
